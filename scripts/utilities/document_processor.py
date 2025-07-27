#!/usr/bin/env python3
"""
Document Processing System
Extracts and converts content from various file types to markdown format:
- PDFs: Uses pdftotext or OCR tools for text extraction
- DOCX/PPTX: Extract text using python-docx and python-pptx
- Images: Apply OCR using tesseract if needed
- Videos: Extract audio and transcribe using whisper
- Convert all content to markdown format for consistency
"""

import os
import sys
import subprocess
import logging
import tempfile
import shutil
from pathlib import Path
from typing import Dict, List, Optional, Union, Tuple
from datetime import datetime
import mimetypes
import hashlib
import json
import re

# Document processing imports
try:
    import pytesseract
    from PIL import Image
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    logging.warning("Tesseract OCR not available. Image text extraction will be disabled.")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available. DOCX processing will be disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PPTX processing will be disabled.")

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    logging.warning("PyPDF2 not available. PDF processing will use pdftotext only.")

try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False
    logging.warning("Whisper not available. Video transcription will be disabled.")


class DocumentProcessor:
    """
    Comprehensive document processing system that can handle various file types
    and convert them to markdown format for AI readability.
    """
    
    SUPPORTED_EXTENSIONS = {
        'pdf': ['pdf'],
        'docx': ['docx', 'doc'],
        'pptx': ['pptx', 'ppt'],
        'image': ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp'],
        'video': ['mp4', 'avi', 'mkv', 'webm', 'mov', 'flv', 'm4v'],
        'text': ['txt', 'md', 'rtf', 'csv'],
        'other': ['html', 'htm', 'xml', 'json']
    }
    
    def __init__(self, output_dir: str = "processed_documents", temp_dir: Optional[str] = None):
        """Initialize the document processor."""
        self.output_dir = Path(output_dir)
        self.temp_dir = Path(temp_dir) if temp_dir else Path(tempfile.gettempdir()) / "doc_processor"
        
        # Create directories
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        
        # Setup logging
        self.setup_logging()
        
        # Initialize whisper model if available
        self.whisper_model = None
        if WHISPER_AVAILABLE:
            try:
                self.whisper_model = whisper.load_model("base")
                self.logger.info("Whisper model loaded successfully")
            except Exception as e:
                self.logger.warning(f"Could not load Whisper model: {e}")
        
        # Processing statistics
        self.stats = {
            'processed': 0,
            'success': 0,
            'failed': 0,
            'by_type': {},
            'start_time': datetime.now()
        }
        
    def setup_logging(self):
        """Setup logging configuration."""
        log_file = self.output_dir / "document_processing.log"
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(log_file),
                logging.StreamHandler(sys.stdout)
            ]
        )
        self.logger = logging.getLogger(__name__)
        
    def get_file_type(self, file_path: Union[str, Path]) -> str:
        """Determine the file type based on extension."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower().lstrip('.')
        
        for file_type, extensions in self.SUPPORTED_EXTENSIONS.items():
            if extension in extensions:
                return file_type
        return 'unknown'
    
    def process_file(self, file_path: Union[str, Path], output_name: Optional[str] = None) -> Dict:
        """
        Process a single file and convert to markdown.
        
        Args:
            file_path: Path to the file to process
            output_name: Optional custom output filename
            
        Returns:
            Dictionary with processing results
        """
        file_path = Path(file_path)
        if not file_path.exists():
            return {'success': False, 'error': 'File not found'}
        
        self.stats['processed'] += 1
        file_type = self.get_file_type(file_path)
        
        # Update type statistics
        self.stats['by_type'][file_type] = self.stats['by_type'].get(file_type, 0) + 1
        
        self.logger.info(f"Processing {file_type} file: {file_path}")
        
        try:
            # Generate output filename
            if not output_name:
                output_name = f"{file_path.stem}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            
            markdown_file = self.output_dir / f"{output_name}.md"
            
            # Process based on file type
            if file_type == 'pdf':
                result = self.process_pdf(file_path, markdown_file)
            elif file_type == 'docx':
                result = self.process_docx(file_path, markdown_file)
            elif file_type == 'pptx':
                result = self.process_pptx(file_path, markdown_file)
            elif file_type == 'image':
                result = self.process_image(file_path, markdown_file)
            elif file_type == 'video':
                result = self.process_video(file_path, markdown_file)
            elif file_type == 'text':
                result = self.process_text_file(file_path, markdown_file)
            elif file_type == 'other':
                result = self.process_other_file(file_path, markdown_file)
            else:
                result = {'success': False, 'error': f'Unsupported file type: {file_type}'}
            
            if result.get('success'):
                self.stats['success'] += 1
                self.logger.info(f"Successfully processed: {file_path}")
            else:
                self.stats['failed'] += 1
                self.logger.error(f"Failed to process {file_path}: {result.get('error', 'Unknown error')}")
            
            return result
            
        except Exception as e:
            self.stats['failed'] += 1
            error_msg = f"Unexpected error processing {file_path}: {str(e)}"
            self.logger.error(error_msg)
            return {'success': False, 'error': error_msg}
    
    def process_pdf(self, file_path: Path, output_file: Path) -> Dict:
        """Process PDF files using pdftotext or OCR."""
        try:
            text_content = ""
            metadata = {}
            
            # First try pdftotext
            try:
                result = subprocess.run(['pdftotext', str(file_path), '-'], 
                                      capture_output=True, text=True, check=True)
                text_content = result.stdout
                method = "pdftotext"
            except (subprocess.CalledProcessError, FileNotFoundError):
                self.logger.warning(f"pdftotext failed for {file_path}, trying PyPDF2")
                
                # Try PyPDF2
                if PYPDF2_AVAILABLE:
                    try:
                        with open(file_path, 'rb') as file:
                            pdf_reader = PyPDF2.PdfReader(file)
                            text_content = ""
                            for page in pdf_reader.pages:
                                text_content += page.extract_text() + "\n"
                            method = "PyPDF2"
                            metadata['pages'] = len(pdf_reader.pages)
                    except Exception as e:
                        self.logger.warning(f"PyPDF2 failed: {e}")
                        text_content = ""
                
                # If text extraction failed, try OCR
                if not text_content.strip() and TESSERACT_AVAILABLE:
                    self.logger.info(f"Attempting OCR for {file_path}")
                    text_content = self.ocr_pdf(file_path)
                    method = "OCR"
            
            if not text_content.strip():
                return {'success': False, 'error': 'No text could be extracted from PDF'}
            
            # Create markdown content
            markdown_content = self.create_markdown_header(file_path, "PDF Document", method, metadata)
            markdown_content += self.clean_text_content(text_content)
            
            # Save markdown file
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'output_file': str(output_file),
                'method': method,
                'character_count': len(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def ocr_pdf(self, file_path: Path) -> str:
        """Perform OCR on PDF pages."""
        try:
            # Convert PDF to images and OCR each page
            temp_dir = self.temp_dir / f"pdf_ocr_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            temp_dir.mkdir(exist_ok=True)
            
            # Convert PDF to images using pdftoppm
            try:
                subprocess.run(['pdftoppm', '-png', str(file_path), str(temp_dir / 'page')], 
                              check=True, capture_output=True)
            except (subprocess.CalledProcessError, FileNotFoundError):
                # Fallback: try with poppler-utils
                try:
                    subprocess.run(['pdftoppm', '-png', str(file_path), str(temp_dir / 'page')], 
                                  check=True)
                except:
                    return ""
            
            # OCR each image
            text_content = ""
            image_files = sorted(temp_dir.glob('page-*.png'))
            
            for image_file in image_files:
                try:
                    image = Image.open(image_file)
                    page_text = pytesseract.image_to_string(image, lang='eng')
                    text_content += f"\n--- Page {image_file.stem.split('-')[1]} ---\n"
                    text_content += page_text + "\n"
                except Exception as e:
                    self.logger.warning(f"OCR failed for {image_file}: {e}")
            
            # Cleanup temp files
            shutil.rmtree(temp_dir, ignore_errors=True)
            
            return text_content
            
        except Exception as e:
            self.logger.error(f"OCR processing failed: {e}")
            return ""
    
    def process_docx(self, file_path: Path, output_file: Path) -> Dict:
        """Process DOCX files using python-docx."""
        if not DOCX_AVAILABLE:
            return {'success': False, 'error': 'python-docx not available'}
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract text content
            text_content = ""
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract tables
            table_content = ""
            for table in doc.tables:
                table_content += "\n### Table\n\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    table_content += f"| {row_text} |\n"
                table_content += "\n"
            
            # Combine content
            full_content = text_content + table_content
            
            if not full_content.strip():
                return {'success': False, 'error': 'No text content found in DOCX'}
            
            # Create markdown
            markdown_content = self.create_markdown_header(file_path, "Word Document", "python-docx")
            markdown_content += self.clean_text_content(full_content)
            
            # Save markdown file
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'output_file': str(output_file),
                'method': 'python-docx',
                'character_count': len(full_content),
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables)
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def process_pptx(self, file_path: Path, output_file: Path) -> Dict:
        """Process PPTX files using python-pptx."""
        if not PPTX_AVAILABLE:
            return {'success': False, 'error': 'python-pptx not available'}
        
        try:
            prs = Presentation(file_path)
            
            text_content = ""
            slide_count = 0
            
            for i, slide in enumerate(prs.slides, 1):
                slide_count += 1
                text_content += f"\n## Slide {i}\n\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    
                    # Extract text from tables
                    if hasattr(shape, "table"):
                        text_content += "\n### Table\n\n"
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            text_content += f"| {row_text} |\n"
                        text_content += "\n"
                
                text_content += "\n---\n"
            
            if not text_content.strip():
                return {'success': False, 'error': 'No text content found in PPTX'}
            
            # Create markdown
            markdown_content = self.create_markdown_header(file_path, "PowerPoint Presentation", "python-pptx")
            markdown_content += self.clean_text_content(text_content)
            
            # Save markdown file
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'output_file': str(output_file),
                'method': 'python-pptx',
                'character_count': len(text_content),
                'slides': slide_count
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def process_image(self, file_path: Path, output_file: Path) -> Dict:
        """Process image files using OCR."""
        if not TESSERACT_AVAILABLE:
            return {'success': False, 'error': 'Tesseract OCR not available'}
        
        try:
            # Open and process image
            image = Image.open(file_path)
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Extract text using OCR
            text_content = pytesseract.image_to_string(image, lang='eng')
            
            if not text_content.strip():
                return {'success': False, 'error': 'No text found in image'}
            
            # Get image metadata
            metadata = {
                'size': image.size,
                'mode': image.mode,
                'format': image.format
            }
            
            # Create markdown
            markdown_content = self.create_markdown_header(file_path, "Image Document", "Tesseract OCR", metadata)
            markdown_content += self.clean_text_content(text_content)
            
            # Save markdown file
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'output_file': str(output_file),
                'method': 'Tesseract OCR',
                'character_count': len(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def process_video(self, file_path: Path, output_file: Path) -> Dict:
        """Process video files by extracting audio and transcribing."""
        if not WHISPER_AVAILABLE or not self.whisper_model:
            return {'success': False, 'error': 'Whisper not available for video transcription'}
        
        try:
            # Extract audio from video
            audio_file = self.temp_dir / f"{file_path.stem}_audio.wav"
            
            # Use ffmpeg to extract audio
            cmd = [
                'ffmpeg', '-i', str(file_path),
                '-acodec', 'pcm_s16le', '-ac', '1', '-ar', '16000',
                str(audio_file), '-y'
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                return {'success': False, 'error': f'Failed to extract audio: {result.stderr}'}
            
            # Transcribe audio using Whisper
            self.logger.info(f"Transcribing audio from {file_path}...")
            transcript_result = self.whisper_model.transcribe(str(audio_file))
            
            text_content = transcript_result['text']
            
            if not text_content.strip():
                return {'success': False, 'error': 'No speech found in video'}
            
            # Get video metadata
            metadata = {
                'language': transcript_result.get('language', 'unknown'),
                'duration': f"{transcript_result.get('duration', 0):.2f} seconds"
            }
            
            # Create markdown with timestamps if available
            markdown_content = self.create_markdown_header(file_path, "Video Transcript", "Whisper AI", metadata)
            
            # Add detailed transcript with timestamps if available
            if 'segments' in transcript_result:
                markdown_content += "\n## Detailed Transcript\n\n"
                for segment in transcript_result['segments']:
                    start_time = self.format_timestamp(segment.get('start', 0))
                    end_time = self.format_timestamp(segment.get('end', 0))
                    text = segment.get('text', '').strip()
                    markdown_content += f"**[{start_time} - {end_time}]** {text}\n\n"
            else:
                markdown_content += "\n## Transcript\n\n"
                markdown_content += self.clean_text_content(text_content)
            
            # Cleanup audio file
            if audio_file.exists():
                audio_file.unlink()
            
            # Save markdown file
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'output_file': str(output_file),
                'method': 'Whisper AI',
                'character_count': len(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def process_text_file(self, file_path: Path, output_file: Path) -> Dict:
        """Process plain text files."""
        try:
            # Read text content
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
            
            if not text_content.strip():
                return {'success': False, 'error': 'File is empty'}
            
            # Create markdown
            markdown_content = self.create_markdown_header(file_path, "Text Document", "Direct conversion")
            markdown_content += self.clean_text_content(text_content)
            
            # Save markdown file
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'output_file': str(output_file),
                'method': 'Direct conversion',
                'character_count': len(text_content)
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def process_other_file(self, file_path: Path, output_file: Path) -> Dict:
        """Process other file types (HTML, XML, JSON)."""
        try:
            extension = file_path.suffix.lower()
            
            # Read content
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            if not content.strip():
                return {'success': False, 'error': 'File is empty'}
            
            # Process based on type
            if extension in ['.html', '.htm']:
                # Simple HTML to text conversion
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'html.parser')
                text_content = soup.get_text()
                method = "HTML parsing"
            elif extension == '.xml':
                # Basic XML structure preservation
                text_content = content
                method = "XML direct"
            elif extension == '.json':
                # Pretty print JSON
                try:
                    import json
                    data = json.loads(content)
                    text_content = json.dumps(data, indent=2)
                    method = "JSON formatting"
                except:
                    text_content = content
                    method = "JSON direct"
            else:
                text_content = content
                method = "Direct conversion"
            
            # Create markdown
            markdown_content = self.create_markdown_header(file_path, f"{extension.upper()} Document", method)
            if extension == '.json':
                markdown_content += f"\n```json\n{text_content}\n```\n"
            elif extension in ['.xml']:
                markdown_content += f"\n```xml\n{text_content}\n```\n"
            else:
                markdown_content += self.clean_text_content(text_content)
            
            # Save markdown file
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'output_file': str(output_file),
                'method': method,
                'character_count': len(text_content)
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def create_markdown_header(self, file_path: Path, doc_type: str, method: str, metadata: Optional[Dict] = None) -> str:
        """Create a standardized markdown header."""
        header = f"# {file_path.name}\n\n"
        header += f"**Document Type:** {doc_type}\n"
        header += f"**Source File:** {file_path}\n"
        header += f"**Processing Method:** {method}\n"
        header += f"**Processed Date:** {datetime.now().isoformat()}\n"
        header += f"**File Size:** {file_path.stat().st_size} bytes\n"
        
        if metadata:
            header += "\n**Metadata:**\n"
            for key, value in metadata.items():
                header += f"- **{key.title()}:** {value}\n"
        
        header += "\n---\n\n"
        return header
    
    def clean_text_content(self, text: str) -> str:
        """Clean and format text content for markdown."""
        # Remove excessive whitespace
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            cleaned_line = line.strip()
            if cleaned_line:
                cleaned_lines.append(cleaned_line)
            elif cleaned_lines and cleaned_lines[-1]:  # Keep single empty lines
                cleaned_lines.append('')
        
        return '\n'.join(cleaned_lines)
    
    def format_timestamp(self, seconds: float) -> str:
        """Format seconds to MM:SS or HH:MM:SS format."""
        hours = int(seconds // 3600)
        minutes = int((seconds % 3600) // 60)
        seconds = int(seconds % 60)
        
        if hours > 0:
            return f"{hours:02d}:{minutes:02d}:{seconds:02d}"
        else:
            return f"{minutes:02d}:{seconds:02d}"
    
    def process_directory(self, directory_path: Union[str, Path], recursive: bool = True) -> Dict:
        """
        Process all supported files in a directory.
        
        Args:
            directory_path: Path to directory to process
            recursive: Whether to process subdirectories
            
        Returns:
            Dictionary with processing summary
        """
        directory_path = Path(directory_path)
        if not directory_path.exists() or not directory_path.is_dir():
            return {'success': False, 'error': 'Directory not found'}
        
        self.logger.info(f"Processing directory: {directory_path}")
        
        # Find all supported files
        all_files = []
        for file_type, extensions in self.SUPPORTED_EXTENSIONS.items():
            for ext in extensions:
                if recursive:
                    all_files.extend(directory_path.rglob(f"*.{ext}"))
                else:
                    all_files.extend(directory_path.glob(f"*.{ext}"))
        
        results = []
        for file_path in all_files:
            if file_path.is_file():
                result = self.process_file(file_path)
                results.append({
                    'file': str(file_path),
                    'result': result
                })
        
        return {
            'success': True,
            'directory': str(directory_path),
            'files_processed': len(results),
            'results': results,
            'statistics': self.get_statistics()
        }
    
    def get_statistics(self) -> Dict:
        """Get processing statistics."""
        elapsed_time = datetime.now() - self.stats['start_time']
        
        return {
            'total_processed': self.stats['processed'],
            'successful': self.stats['success'],
            'failed': self.stats['failed'],
            'success_rate': f"{(self.stats['success'] / max(self.stats['processed'], 1)) * 100:.1f}%",
            'by_type': self.stats['by_type'],
            'elapsed_time': str(elapsed_time),
            'processing_rate': f"{self.stats['processed'] / max(elapsed_time.total_seconds(), 1):.2f} files/second"
        }
    
    def generate_report(self) -> str:
        """Generate a processing report."""
        stats = self.get_statistics()
        
        report = f"""# Document Processing Report

**Generated:** {datetime.now().isoformat()}
**Output Directory:** {self.output_dir}

## Summary Statistics

- **Total Files Processed:** {stats['total_processed']}
- **Successful:** {stats['successful']}
- **Failed:** {stats['failed']}
- **Success Rate:** {stats['success_rate']}
- **Processing Time:** {stats['elapsed_time']}
- **Processing Rate:** {stats['processing_rate']}

## Files by Type

"""
        
        for file_type, count in stats['by_type'].items():
            report += f"- **{file_type.title()}:** {count} files\n"
        
        report += f"""

## System Capabilities

- **PDF Processing:** {'✓' if True else '✗'} (pdftotext + PyPDF2 + OCR)
- **DOCX Processing:** {'✓' if DOCX_AVAILABLE else '✗'} (python-docx)
- **PPTX Processing:** {'✓' if PPTX_AVAILABLE else '✗'} (python-pptx)
- **Image OCR:** {'✓' if TESSERACT_AVAILABLE else '✗'} (Tesseract)
- **Video Transcription:** {'✓' if WHISPER_AVAILABLE else '✗'} (Whisper AI)

## Output Files

All processed documents are converted to markdown format and saved in:
`{self.output_dir}`

Each markdown file includes:
- Source file information
- Processing method used
- Metadata (when available)
- Extracted text content
- Timestamps (for video transcripts)

"""
        
        return report


def main():
    """Command-line interface for document processing."""
    import argparse
    
    parser = argparse.ArgumentParser(description="Process documents and convert to markdown")
    parser.add_argument("input", help="Input file or directory path")
    parser.add_argument("-o", "--output", default="processed_documents", 
                       help="Output directory (default: processed_documents)")
    parser.add_argument("-r", "--recursive", action="store_true",
                       help="Process directories recursively")
    parser.add_argument("--report", action="store_true",
                       help="Generate processing report")
    
    args = parser.parse_args()
    
    processor = DocumentProcessor(output_dir=args.output)
    
    input_path = Path(args.input)
    
    if input_path.is_file():
        result = processor.process_file(input_path)
        print(f"Processing result: {result}")
    elif input_path.is_dir():
        result = processor.process_directory(input_path, recursive=args.recursive)
        print(f"Processed {result.get('files_processed', 0)} files")
        print(f"Statistics: {result.get('statistics', {})}")
    else:
        print(f"Error: {args.input} is not a valid file or directory")
        return 1
    
    if args.report:
        report = processor.generate_report()
        report_file = processor.output_dir / "processing_report.md"
        with open(report_file, 'w', encoding='utf-8') as f:
            f.write(report)
        print(f"Report saved to: {report_file}")
    
    return 0


if __name__ == "__main__":
    sys.exit(main())
